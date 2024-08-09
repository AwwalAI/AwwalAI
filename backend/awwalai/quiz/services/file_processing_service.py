import os
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
import csv

class FileProcessingService:
    
    def process_file(self, file):
        file_extension = os.path.splitext(file.name)[1].lower()

        if file_extension == '.pdf':
            return self.extract_pdf_content(file)
        elif file_extension in ['.doc', '.docx']:
            return self.extract_docx_content(file)
        elif file_extension in ['.xls', '.xlsx']:
            return self.extract_excel_content(file)
        elif file_extension in ['.ppt', '.pptx']:
            return self.extract_pptx_content(file)
        elif file_extension == '.txt':
            return self.extract_text_content(file)
        elif file_extension == '.csv':
            return self.extract_csv_content(file)
        else:
            raise ValueError("Unsupported file type")

    def extract_pdf_content(self, file):
        content = ""
        pdf_reader = PdfReader(file)
        for page in pdf_reader.pages:
            content += page.extract_text()
        return content

    def extract_docx_content(self, file):
        doc = DocxDocument(file)
        return "\n".join([para.text for para in doc.paragraphs])

    def extract_excel_content(self, file):
        wb = load_workbook(file)
        content = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                content.append(" ".join([str(cell) for cell in row if cell]))
        return "\n".join(content)

    def extract_pptx_content(self, file):
        prs = Presentation(file)
        content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content.append(shape.text)
        return "\n".join(content)

    def extract_text_content(self, file):
        return file.read().decode('utf-8')

    def extract_csv_content(self, file):
        content = []
        csv_reader = csv.reader(file.read().decode('utf-8').splitlines())
        for row in csv_reader:
            content.append(" ".join(row))
        return "\n".join(content)
